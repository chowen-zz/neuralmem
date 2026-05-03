"""Office document extractor — Word, Excel, PowerPoint.

Uses the following libraries:

* ``python-docx`` for ``.docx`` files.
* ``openpyxl`` for ``.xlsx`` files.
* ``python-pptx`` for ``.pptx`` files.

If a library is missing the corresponding format falls back to a raw text
decoding (which usually produces garbage for true binary Office files, but
may work for text-based legacy formats such as ``.csv`` exported as ``.xlsx``).
"""
from __future__ import annotations

import logging
import os
import pathlib
from typing import TYPE_CHECKING

from neuralmem.multimodal.base import BaseMultimodalExtractor, ExtractionResult

if TYPE_CHECKING:
    from collections.abc import Sequence

_logger = logging.getLogger(__name__)

# Graceful dependency checks ------------------------------------------------
try:
    import docx  # type: ignore[import-untyped]

    _HAS_DOCX = True
except ImportError:
    docx = None  # type: ignore[misc]
    _HAS_DOCX = False

try:
    import openpyxl  # type: ignore[import-untyped]

    _HAS_OPENPYXL = True
except ImportError:
    openpyxl = None  # type: ignore[misc]
    _HAS_OPENPYXL = False

try:
    import pptx  # type: ignore[import-untyped]

    _HAS_PPTX = True
except ImportError:
    pptx = None  # type: ignore[misc]
    _HAS_PPTX = False


class OfficeExtractor(BaseMultimodalExtractor):
    """Extract text from Microsoft Office documents.

    Args:
        max_chunk_size: Maximum characters per returned chunk.
    """

    def __init__(self, max_chunk_size: int = 4096) -> None:
        super().__init__(max_chunk_size)

        missing = []
        if not _HAS_DOCX:
            missing.append("python-docx")
        if not _HAS_OPENPYXL:
            missing.append("openpyxl")
        if not _HAS_PPTX:
            missing.append("python-pptx")
        if missing:
            _logger.warning(
                "Missing office libraries: %s. "
                "OfficeExtractor will fall back to raw text for those formats.",
                ", ".join(missing),
            )

    def extract(self, path_or_bytes: str | os.PathLike[str] | bytes) -> ExtractionResult:
        """Extract text from an Office document.

        The file extension is used to decide which sub-extractor to invoke:
        ``.docx`` → Word, ``.xlsx`` → Excel, ``.pptx`` → PowerPoint.

        Args:
            path_or_bytes: Path to an office file or raw bytes.

        Returns:
            List of text chunks.
        """
        raw = self._read_bytes(path_or_bytes)
        if not raw:
            return []

        # Determine format from extension when possible
        ext = ""
        if isinstance(path_or_bytes, (str, os.PathLike)):
            ext = pathlib.Path(path_or_bytes).suffix.lower()

        if ext == ".docx":
            return self._extract_docx(raw)
        if ext == ".xlsx":
            return self._extract_xlsx(raw)
        if ext == ".pptx":
            return self._extract_pptx(raw)

        # Unknown / no extension — try each format in turn
        for method in (self._extract_docx, self._extract_xlsx, self._extract_pptx):
            result = method(raw)
            if result:
                return result

        return self._fallback(raw)

    # ------------------------------------------------------------------ #
    # Word (.docx)
    # ------------------------------------------------------------------ #

    def _extract_docx(self, data: bytes) -> ExtractionResult:
        """Extract text from a Word document."""
        if not _HAS_DOCX or docx is None:
            return self._fallback(data)

        try:
            import io

            doc = docx.Document(io.BytesIO(data))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n\n".join(paragraphs)
            return self._chunk_text(text) if text else []
        except Exception as exc:
            _logger.warning("python-docx failed: %s", exc)
            return self._fallback(data)

    # ------------------------------------------------------------------ #
    # Excel (.xlsx)
    # ------------------------------------------------------------------ #

    def _extract_xlsx(self, data: bytes) -> ExtractionResult:
        """Extract text from an Excel workbook."""
        if not _HAS_OPENPYXL or openpyxl is None:
            return self._fallback(data)

        try:
            import io

            wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
            texts: list[str] = []
            for sheet in wb.worksheets:
                rows: list[str] = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        rows.append(row_text)
                if rows:
                    texts.append(f"--- Sheet: {sheet.title} ---\n" + "\n".join(rows))
            return self._chunk_text("\n\n".join(texts)) if texts else []
        except Exception as exc:
            _logger.warning("openpyxl failed: %s", exc)
            return self._fallback(data)

    # ------------------------------------------------------------------ #
    # PowerPoint (.pptx)
    # ------------------------------------------------------------------ #

    def _extract_pptx(self, data: bytes) -> ExtractionResult:
        """Extract text from a PowerPoint presentation."""
        if not _HAS_PPTX or pptx is None:
            return self._fallback(data)

        try:
            import io

            prs = pptx.Presentation(io.BytesIO(data))
            texts: list[str] = []
            for i, slide in enumerate(prs.slides, start=1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    texts.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))
            return self._chunk_text("\n\n".join(texts)) if texts else []
        except Exception as exc:
            _logger.warning("python-pptx failed: %s", exc)
            return self._fallback(data)

    # ------------------------------------------------------------------ #
    # Fallback
    # ------------------------------------------------------------------ #

    def _fallback(self, data: bytes) -> ExtractionResult:
        """Decode raw bytes as UTF-8 with replacement characters."""
        text = data.decode("utf-8", errors="replace")
        return self._chunk_text(text)
